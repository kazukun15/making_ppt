import streamlit as st
import os, time, re, json, requests, tempfile, asyncio
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from langchain.prompts import PromptTemplate
from langchain import LLMChain
from langchain_google_genai import ChatGoogleGenerativeAI

# === APIキーの取得（st.secretsから） ===
GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
TAVILY_API_KEY = st.secrets["TAVILY_API_KEY"]

# === 各エージェントの初期化 (すべて gemini-2.0-flash モデル使用) ===
llm_content = ChatGoogleGenerativeAI(
    model="gemini-2.0-flash",
    temperature=0.3,
    google_api_key=GOOGLE_API_KEY
)
llm_research = ChatGoogleGenerativeAI(
    model="gemini-2.0-flash",
    temperature=0.3,
    google_api_key=GOOGLE_API_KEY
)
llm_enhancement = ChatGoogleGenerativeAI(
    model="gemini-2.0-flash",
    temperature=0.3,
    google_api_key=GOOGLE_API_KEY
)
llm_visualization = ChatGoogleGenerativeAI(
    model="gemini-2.0-flash",
    temperature=0.3,
    google_api_key=GOOGLE_API_KEY
)

# グローバル変数（出典情報など）
source_url_mapping = {}
unique_references = []

# --- 共通ヘルパー関数 ---
def tavily_search(query, api_key, limit=10):
    url = "https://api.tavily.com/search"
    payload = {
        "query": query,
        "api_key": api_key,
        "limit": limit,
        "search_depth": "advanced",
        "topic": "general",
        "include_images": False,
        "include_answer": "advanced",
        "include_raw_content": False
    }
    response = requests.post(url, json=payload)
    if response.status_code == 200:
        return response.json()
    else:
        st.error(f"Tavily API error: {response.status_code} {response.text}")
        return None

def format_search_results(results):
    if not results or "results" not in results:
        return ""
    formatted = ""
    for result in results["results"]:
        title = result.get("title", "")
        description = result.get("description", "")
        url = result.get("url", "")
        formatted += f"Title: {title}\nDescription: {description}\nURL: {url}\n\n"
    return formatted.strip()

def collect_references(results):
    if not results or "results" not in results:
        return []
    references = []
    for result in results["results"]:
        if "title" in result and "url" in result:
            title = result.get("title", "").strip()
            url = result.get("url", "").strip()
            if title and url:
                references.append({"title": title, "url": url})
    return references

def invoke_with_retry(chain, inputs, max_retries=7):
    retries = 0
    backoff = 5
    while True:
        try:
            return chain.invoke(inputs)
        except Exception as e:
            st.warning(f"リクエストエラー: {e}. {backoff}秒後に再試行します...")
            time.sleep(backoff)
            retries += 1
            backoff *= 2
            if retries >= max_retries:
                raise e

def extract_json(response_text):
    st.write(f"JSON抽出開始（文字数: {len(response_text)}）")
    if not response_text or not response_text.strip():
        return '{"slides":[]}'

    code_block_patterns = [
        r'```(?:json)?\s*([\s\S]*?)\s*```',
        r'```\s*json\s*\n([\s\S]*?)\n\s*```'
    ]
    for pattern in code_block_patterns:
        code_blocks = re.findall(pattern, response_text, re.DOTALL)
        for block in code_blocks:
            try:
                parsed = json.loads(block.strip())
                if "slides" in parsed:
                    return block.strip()
            except json.JSONDecodeError:
                continue

    json_pattern = r'({[\s\S]*})'
    matches = re.findall(json_pattern, response_text)
    for match in matches:
        try:
            parsed = json.loads(match)
            if "slides" in parsed:
                return match
        except:
            pass
    return '{"slides":[]}'

def extract_table_json(response_text):
    if not response_text or not response_text.strip():
        return '{"table_title":"データ表","headers":["項目"],"rows":[["データなし"]]}'
    code_block_patterns = [
        r'```(?:json)?\s*([\s\S]*?)\s*```',
        r'```\s*json\s*\n([\s\S]*?)\n\s*```'
    ]
    for pattern in code_block_patterns:
        code_blocks = re.findall(pattern, response_text, re.DOTALL)
        for block in code_blocks:
            try:
                parsed = json.loads(block.strip())
                if "table_title" in parsed:
                    return block.strip()
            except json.JSONDecodeError:
                continue
    json_pattern = r'({[\s\S]*})'
    matches = re.findall(json_pattern, response_text)
    for match in matches:
        try:
            parsed = json.loads(match)
            if "table_title" in parsed:
                return match
        except:
            pass
    return '{"table_title":"データ表","headers":["項目1","項目2"],"rows":[["データ1","データ2"],["データ3","データ4"]]}'

def extract_url_from_source(source_text):
    url_match = re.search(r'https?://[^\s)"]+', source_text)
    if url_match:
        return url_match.group(0)
    return None

def add_source_textbox(slide, source_text):
    if not source_text:
        return
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    url = extract_url_from_source(source_text)
    if url:
        parts = re.split(r'(https?://[^\s)"]+)', source_text, 1)
        p = tf.paragraphs[0]
        p.text = parts[0].strip()
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.name = "宮澄乃"
        if len(parts) > 1:
            run = p.add_run()
            run.text = parts[1]
            run.hyperlink.address = url
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.font.underline = True
            run.font.name = "宮澄乃"
        if len(parts) > 2:
            run = p.add_run()
            run.text = parts[2]
            run.font.italic = True
            run.font.name = "宮澄乃"
    else:
        p = tf.paragraphs[0]
        p.text = source_text
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.name = "宮澄乃"
    return txBox

def set_bullet_style(paragraph, enable_bullet=True, level=0):
    paragraph.level = level if enable_bullet else 0

# === 1. Research Agent ===
def research_agent(user_prompt):
    optimize_prompt_template = """
    あなたはウェブ検索クエリの最適化に優れたAIです。
    以下のテーマと指示に基づき、tavilyで検索するための最適な検索クエリを生成してください。生成するクエリは、もとの検索クエリ1個と追加クエリ2個の合計3個である必要があります。
    クエリは異なる視点をカバーし、統計データや定量分析に有用なキーワードのみを返してください。
    必ず以下のJSON形式で返してください：
    {{
      "queries": [
        "元の検索クエリ",
        "第一の検索クエリ",
        "第二の検索クエリ"
      ]
    }}
    ユーザーの入力: {user_input}
    """
    prompt = PromptTemplate(input_variables=["user_input"], template=optimize_prompt_template)
    chain = LLMChain(llm=llm_research, prompt=prompt)
    optimized_response = invoke_with_retry(chain, {"user_input": user_prompt})
    try:
        json_str = extract_json(optimized_response["text"])
        json_data = json.loads(json_str)
        optimized_queries = json_data["queries"]
    except Exception as e:
        st.warning(f"Research JSON解析エラー: {e}")
        optimized_queries = [f"{user_prompt} 概要", f"{user_prompt} 詳細", f"{user_prompt} 出典"]
    data_focused = []
    for query in optimized_queries[:1]:
        data_focused.append(f"{query} 統計 データ")
        data_focused.append(f"{query} 出典 公式")
    all_queries = list(dict.fromkeys([user_prompt] + optimized_queries + data_focused))
    all_web_results = ""
    all_references = []
    for query in all_queries:
        results = tavily_search(query, TAVILY_API_KEY)
        web_text = format_search_results(results)
        all_web_results += f"【検索クエリ: {query}】\n{web_text}\n\n"
        refs = collect_references(results)
        all_references.extend(refs)
    # 重複削除
    seen = set()
    unique_refs = []
    for ref in all_references:
        if ref["url"] not in seen:
            unique_refs.append(ref)
            seen.add(ref["url"])
    return {
        "optimized_queries": optimized_queries,
        "web_results": all_web_results,
        "references": unique_refs
    }

# === 2. Content Planning Agent ===
def content_planning_agent(user_prompt, web_results, total_slide_count):
    presentation_template = f"""
    あなたは高度なプレゼンテーション構成作成AIです。以下のルールを厳守してください：
    1. 各スライドのタイトルは20ポイント、本文（箇条書き・説明文）は18ポイント、出典情報は14ポイントとすること。
    2. 各スライドは必ず6〜8個の箇条書き項目（各60〜80文字程度）と400〜500文字の詳細な説明文、出典情報を含むこと。
    3. 全体で約{total_slide_count}枚程度のスライドを生成し、最後に必ず「まとめ」スライドを1枚含むこと。
    4. 出力は必ず下記のJSON形式のみで返すこと：
    {{
      "slides": [
        {{
          "title": "タイトルスライド",
          "subtitle": "サブタイトル"
        }},
        {{
          "title": "スライド2のタイトル",
          "bullets": [
            "1つ目の箇条書き",
            "2つ目の箇条書き",
            "3つ目の箇条書き",
            "4つ目の箇条書き",
            "5つ目の箇条書き",
            "6つ目の箇条書き"
          ],
          "description": "説明文...",
          "source": "出典情報"
        }}
      ]
    }}
    ユーザーのテーマ：{{prompt}}
    参考Web検索結果：
    {{web_results}}
    """
    prompt = PromptTemplate(input_variables=["prompt", "web_results"], template=presentation_template)
    chain = LLMChain(llm=llm_content, prompt=prompt)
    result = invoke_with_retry(chain, {"prompt": user_prompt, "web_results": web_results})
    json_str = extract_json(result["text"])
    try:
        slides_data = json.loads(json_str)
    except Exception as e:
        st.warning(f"Content Planning JSON解析エラー: {e}")
        slides_data = {"slides": []}
    return slides_data

# === 3. Enhancement Agent ===
def enhancement_agent(existing_slides_data, user_prompt, web_results, total_slide_count):
    enhancement_prompt = f"""
    以下のJSONデータはプレゼンテーションのスライド構成です。いくつかのスライドで本文や出典情報が不足しています。
    既存情報を保持しつつ、以下のルールに従い不足部分を補完してください：
    1. 全体で約{total_slide_count}枚のスライドと「まとめ」スライドを含むこと。
    2. 各スライドは6〜8個の箇条書き項目（各60〜80文字）と400〜500文字の詳細な説明文、出典情報を含むこと。
    3. 出力は余計なテキストを含まず、必ずJSON形式で返すこと。
    ユーザーのテーマ：{{prompt}}
    参考Web検索結果：
    {{web_results}}
    既存のJSONデータ：
    {{existing_json}}
    """
    prompt = PromptTemplate(input_variables=["prompt", "web_results", "existing_json"], template=enhancement_prompt)
    chain = LLMChain(llm=llm_enhancement, prompt=prompt)
    result = invoke_with_retry(chain, {
        "prompt": user_prompt,
        "web_results": web_results,
        "existing_json": json.dumps(existing_slides_data)
    })
    json_str = extract_json(result["text"])
    try:
        enhanced_data = json.loads(json_str)
    except Exception as e:
        st.warning(f"Enhancement JSON解析エラー: {e}")
        enhanced_data = existing_slides_data
    return enhanced_data

# === 4. Quality Check Agent ===
def quality_check_agent(slides_data):
    slides = slides_data.get("slides", [])
    if not any("まとめ" in slide.get("title", "") for slide in slides):
        st.info("まとめスライドが見つからなかったため、デフォルトまとめスライドを追加します。")
        summary_slide = {
            "title": "まとめ",
            "bullets": ["本文なし"],
            "description": "本文なし",
            "source": "本文なし"
        }
        slides.append(summary_slide)
    for slide in slides:
        if not slide.get("bullets") and not slide.get("description"):
            slide["bullets"] = ["本文なし"]
            slide["description"] = "本文なし"
    slides_data["slides"] = slides
    return slides_data

# === 5. Design Agent (PPT生成、フォントは「宮澄乃」、ワンシート内に収まるようテキスト折り返し調整) ===
def design_agent(slides_data):
    prs = Presentation()
    # タイトルスライド
    title_slide_layout = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(title_slide_layout)
    slide0.shapes.title.text = slides_data["slides"][0].get("title", "タイトルなし")
    for paragraph in slide0.shapes.title.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "宮澄乃"
    if "subtitle" in slides_data["slides"][0]:
        slide0.placeholders[1].text = slides_data["slides"][0]["subtitle"]
        for paragraph in slide0.placeholders[1].text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.name = "宮澄乃"
    # 他スライド
    for slide_info in slides_data["slides"][1:]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_info.get("title", "タイトルなし")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.name = "宮澄乃"
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)
        body_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = body_box.text_frame
        text_frame.word_wrap = True
        bullets = slide_info.get("bullets", [])
        if bullets and bullets != ["本文なし"]:
            text_frame.text = bullets[0]
            text_frame.paragraphs[0].font.size = Pt(18)
            text_frame.paragraphs[0].font.name = "宮澄乃"
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.name = "宮澄乃"
        else:
            text_frame.text = "本文なし"
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = "宮澄乃"
        if slide_info.get("description", "").strip() != "":
            p = text_frame.add_paragraph()
            p.text = slide_info["description"]
            p.font.size = Pt(18)
            p.font.name = "宮澄乃"
        if slide_info.get("source", ""):
            add_source_textbox(slide, slide_info["source"])
    return prs

# === 6. Visualization Agent (テーブル生成、セル内折り返し＆ワンシート内に収まるよう調整) ===
def visualization_agent(prs, slides_data, web_results):
    table_prompt_template = """
    あなたは正確なデータ視覚化AIです。以下の情報源をもとに、1ページに収まるサイズの表を生成してください。
    情報はスライド内容と信頼性の高いWeb検索結果から抽出してください。
    出力は以下のJSON形式で返してください：
    {{
      "table_title": "表のタイトル",
      "headers": ["列1", "列2", "列3"],
      "rows": [
        ["行1列1", "行1列2", "行1列3"],
        ["行2列1", "行2列2", "行2列3"]
      ],
      "source": "データの出典情報"
    }}
    スライド内容：
    {slide_content}
    Web検索結果：
    {web_results}
    """
    prompt = PromptTemplate(input_variables=["slide_content", "web_results"], template=table_prompt_template)
    chain = LLMChain(llm=llm_visualization, prompt=prompt)
    try:
        blank_slide_layout = prs.slide_layouts[6]
    except IndexError:
        blank_slide_layout = prs.slide_layouts[5]
    for slide_info in slides_data.get("slides", [])[1:]:
        slide_content_parts = []
        if "title" in slide_info:
            slide_content_parts.append(f"タイトル: {slide_info['title']}")
        if "bullets" in slide_info:
            slide_content_parts.append("箇条書き: " + " | ".join(slide_info["bullets"]))
        if "description" in slide_info:
            slide_content_parts.append(f"説明: {slide_info['description']}")
        slide_content = "\n".join(slide_content_parts)
        filtered_web_results = web_results[:3000]
        try:
            table_response = invoke_with_retry(chain, {"slide_content": slide_content, "web_results": filtered_web_results})
            table_json_str = extract_table_json(table_response["text"])
            table_data = json.loads(table_json_str)
            table_slide = prs.slides.add_slide(blank_slide_layout)
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_box = table_slide.shapes.add_textbox(left, top, width, height)
            title_box.text = f"{slide_info.get('title','')} - {table_data.get('table_title','表')}"
            for paragraph in title_box.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.name = "宮澄乃"
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            row_count = len(rows) + 1
            col_count = len(headers)
            table_left = Inches(0.5)
            table_top = Inches(1.5)
            table_width = Inches(9)
            table_height = Inches(4)
            pptx_table = table_slide.shapes.add_table(row_count, col_count, table_left, table_top, table_width, table_height).table
            for col_idx, header in enumerate(headers):
                cell = pptx_table.cell(0, col_idx)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(18)
                    paragraph.font.name = "宮澄乃"
                    paragraph.alignment = PP_ALIGN.CENTER
                    cell.text_frame.word_wrap = True
            for row_idx, row_data in enumerate(rows, start=1):
                for col_idx, cell_text in enumerate(row_data):
                    cell = pptx_table.cell(row_idx, col_idx)
                    cell.text = cell_text
                    cell.text_frame.word_wrap = True
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.name = "宮澄乃"
                        if cell_text.replace(',', '').replace('.', '').replace('-', '').isdigit():
                            paragraph.alignment = PP_ALIGN.RIGHT
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT
            if col_count == 3:
                pptx_table.columns[0].width = Inches(2.5)
                pptx_table.columns[1].width = Inches(3.0)
                pptx_table.columns[2].width = Inches(3.5)
            if slide_info.get("source", ""):
                add_source_textbox(prs.slides[-1], slide_info["source"])
        except Exception as e:
            st.warning(f"テーブル生成エラー: {e}")
            continue
    if unique_references:
        slide_layout = prs.slide_layouts[1]
        ref_slide = prs.slides.add_slide(slide_layout)
        ref_slide.shapes.title.text = "出典情報"
        for paragraph in ref_slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.name = "宮澄乃"
        body_box = ref_slide.shapes.placeholders[1]
        body_box.text = ""
        for ref in unique_references[:15]:
            p = body_box.text_frame.add_paragraph()
            p.text = ref["title"]
            p.font.size = Pt(20)
            p.font.name = "宮澄乃"
            p.level = 0
            run = p.runs[0]
            run.hyperlink.address = ref["url"]
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.font.underline = True
    return prs

# === 7. オーケストレーション（非同期処理でバックグラウンド実行、同期処理に見せかける） ===
async def main_orchestration_async(user_prompt, save_dir, total_slide_count, progress_bar, status_text):
    status_text.text("Research Agent 処理中...")
    research_output = await asyncio.to_thread(research_agent, user_prompt)
    progress_bar.progress(20)
    
    status_text.text("Content Planning Agent 処理中...")
    slides_data = await asyncio.to_thread(content_planning_agent, user_prompt, research_output["web_results"], total_slide_count)
    progress_bar.progress(40)
    
    status_text.text("Quality Check Agent 処理中...")
    slides_data = await asyncio.to_thread(quality_check_agent, slides_data)
    progress_bar.progress(50)
    
    status_text.text("Enhancement Agent 処理中...")
    slides_data = await asyncio.to_thread(enhancement_agent, slides_data, user_prompt, research_output["web_results"], total_slide_count)
    progress_bar.progress(60)
    
    status_text.text("Design Agent 処理中...")
    prs = await asyncio.to_thread(design_agent, slides_data)
    progress_bar.progress(80)
    
    status_text.text("Visualization Agent 処理中...")
    prs = await asyncio.to_thread(visualization_agent, prs, slides_data, research_output["web_results"])
    progress_bar.progress(95)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_filename = f"presentation_{timestamp}.pptx"
    save_path = os.path.join(save_dir, pptx_filename)
    await asyncio.to_thread(prs.save, save_path)
    progress_bar.progress(100)
    status_text.text("完了！")
    return save_path

def main():
    st.set_page_config(page_title="パワーポイント作成君", layout="wide")
    st.title("パワーポイント作成君")
    st.markdown("### 高品質なプレゼンテーションを自動生成します")
    
    # サイドバーでスライド枚数を調整
    total_slide_count = st.sidebar.slider("生成するスライド枚数", min_value=10, max_value=50, value=25, step=1)
    
    graphviz_code = """
    digraph {
        rankdir=LR;
        UserInput -> ResearchAgent [label="検索クエリ生成"];
        ResearchAgent -> ContentPlanning [label="構成作成"];
        ContentPlanning -> Enhancement [label="補完・強化"];
        Enhancement -> Design [label="デザイン調整"];
        Design -> Visualization [label="データ視覚化"];
        Visualization -> PPTOutput [label="PPT生成"];
    }
    """
    st.graphviz_chart(graphviz_code)
    
    user_prompt = st.text_area("プレゼンテーションのテーマや指示を入力してください：", height=150)
    
    if not user_prompt:
        st.info("テーマや指示を入力してください。")
        return

    save_dir = os.path.join(tempfile.gettempdir(), "generated_ppts")
    os.makedirs(save_dir, exist_ok=True)
    
    if st.button("プレゼンテーション生成"):
        with st.spinner("プレゼンテーションを生成中です..."):
            progress_bar = st.progress(0)
            status_text = st.empty()
            try:
                final_save_path = asyncio.run(main_orchestration_async(user_prompt, save_dir, total_slide_count, progress_bar, status_text))
                with open(final_save_path, "rb") as f:
                    ppt_data = f.read()
                st.success("プレゼンテーション生成完了！")
                st.download_button(
                    label="PPTXファイルをダウンロード",
                    data=ppt_data,
                    file_name=os.path.basename(final_save_path),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"エラーが発生しました: {e}")

if __name__ == '__main__':
    main()
